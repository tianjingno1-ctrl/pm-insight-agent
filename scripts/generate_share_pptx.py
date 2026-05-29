"""Generate PowerPoint for AI learning share presentation."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "ai-learning-share.pptx"

# Brand colors
VIOLET_DARK = RGBColor(0x4C, 0x1D, 0x95)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
FUCHSIA = RGBColor(0xC0, 0x26, 0xD3)
SLATE = RGBColor(0x47, 0x55, 0x69)
SLATE_LIGHT = RGBColor(0x94, 0xA3, 0xB8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF5, 0xF3, 0xFF)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xE1, 0x1D, 0x48)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = SLATE,
    align=PP_ALIGN.LEFT,
    font_name: str = "微软雅黑",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    p.font.color.rgb = color
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items: list[str],
    *,
    font_size: int = 20,
    color: RGBColor = SLATE,
    spacing: int = 10,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.name = "微软雅黑"
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
    return box


def add_title_slide(prs: Presentation, title: str, subtitle: str, meta: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, VIOLET_DARK)

    # Accent bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.12), prs.slide_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = FUCHSIA
    bar.line.fill.background()

    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.8),
        Inches(8.5),
        Inches(1.2),
        title,
        font_size=40,
        bold=True,
        color=WHITE,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(3.0),
        Inches(8.5),
        Inches(0.8),
        subtitle,
        font_size=22,
        color=RGBColor(0xE9, 0xD5, 0xFF),
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(4.5),
        Inches(8.5),
        Inches(1.0),
        meta,
        font_size=16,
        color=SLATE_LIGHT,
    )


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str] | None = None,
    *,
    highlight: str | None = None,
    notes: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    # Top accent
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    header.fill.solid()
    header.fill.fore_color.rgb = BG_LIGHT
    header.line.fill.background()

    add_textbox(
        slide,
        Inches(0.6),
        Inches(0.25),
        Inches(8.8),
        Inches(0.7),
        title,
        font_size=28,
        bold=True,
        color=VIOLET_DARK,
    )

    y = Inches(1.4)
    if highlight:
        add_textbox(
            slide,
            Inches(0.6),
            y,
            Inches(8.8),
            Inches(0.9),
            highlight,
            font_size=22,
            bold=True,
            color=VIOLET,
        )
        y = Inches(2.3)

    if bullets:
        add_bullets(slide, Inches(0.7), y, Inches(8.6), Inches(4.5), bullets, font_size=20)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    *,
    notes: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = BG_LIGHT
    header_bar.line.fill.background()

    add_textbox(
        slide,
        Inches(0.6),
        Inches(0.25),
        Inches(8.8),
        Inches(0.7),
        title,
        font_size=28,
        bold=True,
        color=VIOLET_DARK,
    )

    cols = len(headers)
    table_shape = slide.shapes.add_table(
        len(rows) + 1, cols, Inches(0.5), Inches(1.35), Inches(9.0), Inches(0.5 + 0.55 * (len(rows) + 1))
    )
    table = table_shape.table

    col_w = Inches(9.0 / cols)
    for i in range(cols):
        table.columns[i].width = int(col_w)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = VIOLET
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.name = "微软雅黑"
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = "微软雅黑"
                p.font.color.rgb = SLATE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
    *,
    notes: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    header_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = BG_LIGHT
    header_bar.line.fill.background()

    add_textbox(slide, Inches(0.6), Inches(0.25), Inches(8.8), Inches(0.7), title, font_size=28, bold=True, color=VIOLET_DARK)

    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(4.2), Inches(0.5), left_title, font_size=18, bold=True, color=GREEN)
    add_bullets(slide, Inches(0.5), Inches(1.85), Inches(4.2), Inches(4.5), left_items, font_size=16)

    add_textbox(slide, Inches(5.0), Inches(1.3), Inches(4.2), Inches(0.5), right_title, font_size=18, bold=True, color=RED)
    add_bullets(slide, Inches(5.0), Inches(1.85), Inches(4.2), Inches(4.5), right_items, font_size=16)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def add_closing_slide(prs: Presentation, quote: str, sub: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, VIOLET_DARK)

    add_textbox(
        slide,
        Inches(0.8),
        Inches(2.0),
        Inches(8.4),
        Inches(2.0),
        quote,
        font_size=28,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(4.2),
        Inches(8.4),
        Inches(1.2),
        sub,
        font_size=18,
        color=RGBColor(0xE9, 0xD5, 0xFF),
        align=PP_ALIGN.CENTER,
    )


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "PM Insight Agent",
        "我的 AI 学习实践分享",
        "分享人：[填写姓名]  ·  需求评审会  ·  2026\n从 0 到 1 搭建「AI 专家圆桌」原型",
    )

    add_content_slide(
        prs,
        "今天要带走什么",
        [
            "看见一种可能 — AI 可以做成「多角色协作」，不只是聊天框",
            "理解边界 — 哪些是 Demo，哪些是真实工程能力",
            "可复用方法 — mock 先行、分阶段交付、Cursor 辅助开发",
        ],
        notes="设定预期，避免听众以为要讲大模型原理。",
    )

    add_table_slide(
        prs,
        "产品经理的日常困境",
        ["场景", "痛点"],
        [
            ["需求评审前", "信息散、难提炼"],
            ["多方争论", "缺结构化讨论"],
            ["写 PRD", "重复劳动多"],
            ["学 AI", "看很多、做很少"],
        ],
        notes="用大家熟悉的场景开场，建立共鸣。",
    )

    add_content_slide(
        prs,
        "项目是什么？",
        [],
        highlight="你提一个议题  →  多个 AI 专家角色讨论  →  输出结构化小结",
        notes="30 秒讲清楚产品形态。可插入圆桌 UI 截图。",
    )

    add_table_slide(
        prs,
        "演示界面长什么样",
        ["圆桌区", "发言记录", "决策小结"],
        [
            ["短气泡显示互动", "全部发言汇总流式", "倾向 / 分歧 / 下一步"],
            ["四专家围坐动画", "自动滚到底部", "三列卡片展示"],
        ],
        notes="演示时指三个区域给观众看。",
    )

    add_content_slide(
        prs,
        "我是怎么一步步做出来的",
        [
            "CLI 报告 → Streamlit → 圆桌 UI → SSE → DeepSeek",
            "当前：发言记录 + 演示级 UI",
            "",
            "关键：mock 先行 → 协议对齐 → 再接真 API",
        ],
        notes="体现 PM 也能分阶段推进。",
    )

    add_content_slide(
        prs,
        "整条链路怎么跑",
        [
            "用户输入议题",
            "      ↓",
            "前端圆桌 UI（展示 + 发言记录流式）",
            "      ↓  SSE 事件流",
            "后端调用 DeepSeek（失败自动降级）",
            "      ↓",
            "输出：多角色发言 + 结构化小结",
        ],
        notes="强调事件协议：换数据源不改 UI。",
    )

    add_two_column_slide(
        prs,
        "当前「是」与「不是」",
        "✅  已经做到",
        [
            "真实 LLM 生成讨论",
            "前后端全链路",
            "API 失败自动降级",
            "可视化圆桌 + 发言记录",
        ],
        "❌  还没做到",
        [
            "真正多 Agent 编排",
            "边生成边播",
            "向量 RAG",
            "生产级部署",
        ],
        notes="主动说边界，建立信任。",
    )

    add_content_slide(
        prs,
        "Live Demo（5 分钟）★",
        [
            "1. 打开 localhost:3000，F11 全屏",
            "2. 输入议题：「需求评审会要不要引入 AI 辅助？」",
            "3. 点击「开始讨论」，等待 10～30 秒",
            "4. 指左侧发言记录 + 圆桌短气泡",
            "5. 指底部本轮决策小结（三列卡片）",
            "",
            "线上会：腾讯会议 / 飞书共享屏幕即可",
        ],
        notes="话术：圆桌看互动，左侧看完整发言，最后看结构化小结。",
    )

    add_table_slide(
        prs,
        "怎么分享给别人看",
        ["方式", "怎么做", "适合"],
        [
            ["投屏演示 ★", "F11 全屏 + 屏幕共享", "评审会、分享会"],
            ["发 Word", "ai-learning-share.docx", "会后留存"],
            ["录屏 MP4", "Win+G 录完整讨论", "异步培训"],
            ["代码仓库", "同事 clone 自己跑", "技术同事"],
        ],
        notes="⚠️ localhost 只有你自己能开，不要发链接给别人。",
    )

    add_content_slide(
        prs,
        "我学到的 5 件事",
        [
            "1. 先 mock，再接真 API",
            "2. 协议比模型更重要",
            "3. 「多 Agent」分演示型 vs 工程型",
            "4. 分阶段交付 + Git 锚点 + 交接文档",
            "5. 演示稳定性 > 技术炫技",
        ],
        notes="分享核心页，每条 30 秒 + 真实例子。",
    )

    add_table_slide(
        prs,
        "AI 能怎么帮需求评审？",
        ["传统", "AI 辅助后"],
        [
            ["人工整理意见", "AI 出讨论草稿"],
            ["现场争论发散", "AI 预演一轮"],
            ["PRD 从空白写", "从材料生成初稿"],
            ["结论靠人记", "固定格式小结"],
        ],
        notes="AI 输出是候选方案，决策仍由人负责。",
    )

    add_content_slide(
        prs,
        "真实踩坑",
        [
            "议题没传给 LLM → 答非所问",
            "两套环境变量 → DeepSeek 接不上",
            "气泡字太多 → 改短标签 + 发言记录",
            "端口僵尸进程 → 演示前必查服务",
        ],
        notes="展示真实学习过程更有说服力。",
    )

    add_table_slide(
        prs,
        "团队落地路径",
        ["阶段", "目标", "投入"],
        [
            ["试点", "CLI 出讨论草稿", "1 人 / 1 周"],
            ["内测", "Streamlit 内网", "1～2 人 / 2 周"],
            ["体验升级", "圆桌 UI 演示", "持续迭代"],
            ["生产化", "鉴权 + RAG + 飞书", "研发排期"],
        ],
    )

    add_content_slide(
        prs,
        "Roadmap",
        [
            "P0：真正多 Agent 编排",
            "P1：边生成边播",
            "P2：会议记录 RAG",
            "P3：飞书 / 钉钉集成",
        ],
    )

    add_content_slide(
        prs,
        "常见问题 Q&A",
        [
            "套壳吗？→ 多角色圆桌 + 事件协议，不是单轮问答",
            "安全吗？→ 本地原型，落地走公司模型网关",
            "非技术能改吗？→ 议题/prompt 可以",
            "成本？→ DeepSeek 一次几分钱量级",
        ],
    )

    add_closing_slide(
        prs,
        "AI 不会取代产品经理，\n但会用 AI 的产品经理会取代不会用的。",
        "目标是跑通「需求 → 多视角讨论 → 结构化输出」的 AI 链路",
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Generated: {OUT}")


if __name__ == "__main__":
    build()
